from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "conference-latex-template" / "IEEE-conference-template-062824" / "figures"
OUT = ROOT / "final-presentation.pptx"

INK = "0F172A"
SLATE = "475569"
MUTED = "64748B"
BLUE = "2563EB"
TEAL = "0F766E"
AMBER = "B45309"
VIOLET = "6D28D9"
GREEN = "15803D"
RED = "B91C1C"
BG = "F8FAFC"
PANEL = "FFFFFF"


def color(hex_value: str) -> RGBColor:
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_bg(slide, hex_value: str = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(hex_value)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 28,
    bold: bool = False,
    font_color: str = INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos Display"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(font_color)


def add_title(slide, title: str, subtitle: str = "", *, speaker: str = "") -> None:
    add_text(slide, title, 0.62, 0.38, 8.5, 0.58, size=30, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.96, 8.7, 0.34, size=13, font_color=MUTED)
    if speaker:
        add_chip(slide, speaker, 10.52, 0.42, 2.2, 0.34, fill="E0F2FE", font_color="075985")


def add_footer(slide, slide_no: int) -> None:
    add_text(slide, "Embodied AI Final Project", 0.62, 7.08, 3.0, 0.22, size=9, font_color=MUTED)
    add_text(slide, f"{slide_no:02d}", 12.5, 7.08, 0.42, 0.22, size=9, font_color=MUTED, align=PP_ALIGN.RIGHT)


def add_chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "E2E8F0",
    font_color: str = INK,
    size: int = 10,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color("CBD5E1")
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color(font_color)


def add_panel(slide, x: float, y: float, w: float, h: float, *, fill: str = PANEL, line: str = "CBD5E1"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(1.0)
    return shape


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, *, size: int = 18, font_color: str = INK) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color(font_color)


def add_metric(slide, label: str, value: str, x: float, y: float, w: float, *, accent: str = BLUE) -> None:
    add_panel(slide, x, y, w, 0.92, fill="FFFFFF", line="CBD5E1")
    add_text(slide, value, x + 0.16, y + 0.13, w - 0.32, 0.34, size=22, bold=True, font_color=accent)
    add_text(slide, label, x + 0.16, y + 0.53, w - 0.32, 0.24, size=10, font_color=MUTED)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    if not path.exists():
        raise FileNotFoundError(path)
    with Image.open(path) as img:
        iw, ih = img.size
    box_w, box_h = Inches(w), Inches(h)
    ratio = min(box_w / iw, box_h / ih)
    new_w, new_h = int(iw * ratio), int(ih * ratio)
    left = Inches(x) + int((box_w - new_w) / 2)
    top = Inches(y) + int((box_h - new_h) / 2)
    image_stream = BytesIO(path.read_bytes())
    slide.shapes.add_picture(image_stream, left, top, width=new_w, height=new_h)


def add_image_panel(slide, path: Path, x: float, y: float, w: float, h: float, *, label: str = "") -> None:
    add_panel(slide, x, y, w, h)
    add_image_fit(slide, path, x + 0.12, y + 0.12, w - 0.24, h - 0.36 if label else h - 0.24)
    if label:
        add_text(slide, label, x + 0.16, y + h - 0.25, w - 0.32, 0.18, size=8, font_color=MUTED, align=PP_ALIGN.CENTER)


def add_arrow(slide, x: float, y: float, w: float, h: float, *, fill: str = INK) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(fill)


def add_simple_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, *, header_fill: str = "DBEAFE") -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c in range(len(rows[0])):
        table.columns[c].width = int(Inches(w) / len(rows[0]))
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = color(header_fill if r == 0 else "FFFFFF")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(10 if r else 11)
                paragraph.font.bold = r == 0
                paragraph.font.color.rgb = color(INK)


def set_notes(slide, speaker: str, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = f"Speaker: {speaker}\n\n{notes.strip()}"


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    return slide


def build() -> None:
    required = [
        "verified_system_pipeline.png",
        "canonical_skeleton_seed.png",
        "real_guided_skeleton_augmentation.png",
        "real_guided_alignment_error.png",
        "real_guided_pose_schedule.png",
        "tcn_temporal_model.png",
        "final_take_timing_budget.png",
        "final_live_candidate_poster.png",
        "three_take_review_poster.png",
        "schunk_rock_yaw45_pitch20.png",
        "schunk_paper_yaw45_pitch20.png",
        "schunk_scissors_yaw45_pitch20.png",
    ]
    for name in required:
        if not (FIGURES / name).exists():
            raise FileNotFoundError(FIGURES / name)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Final Presentation - Real-Guided Skeleton Augmentation for RPS Counterattack"
    prs.core_properties.author = "MinKyu Cho and SeungHwan Kim"

    # 1
    slide = new_slide(prs)
    set_bg(slide, "0B1220")
    add_text(slide, "Real-Guided Skeleton\nAugmentation for\nActuator-Constrained\nRPS Counterattack", 0.7, 0.65, 5.5, 2.35, size=34, bold=True, font_color="F8FAFC")
    add_text(slide, "Embodied Artificial Intelligence Final Project\nMinKyu Cho and SeungHwan Kim · Hanyang University", 0.72, 3.22, 5.35, 0.62, size=14, font_color="CBD5E1")
    add_chip(slide, "Final report + presentation", 0.74, 4.1, 2.25, 0.36, fill="DBEAFE", font_color="1D4ED8")
    add_chip(slide, "15-20 min", 3.14, 4.1, 1.2, 0.36, fill="DCFCE7", font_color="166534")
    add_image_panel(slide, FIGURES / "final_live_candidate_poster.png", 6.45, 0.56, 6.15, 5.85, label="Final live prompt-window counterattack candidate")
    add_footer(slide, 1)
    set_notes(slide, "MinKyu Cho", "안녕하세요. 저희 프로젝트는 가위바위보를 단순히 맞히는 데모가 아니라, 사람 손 동작의 초반 움직임을 보고 로봇 손이 시간 안에 counterattack 할 수 있는지를 검증하는 프로젝트입니다. 오늘 발표에서는 real skeleton data를 어떻게 늘렸고, 그 모델이 live prompt window에서 어떻게 동작했는지를 설명하겠습니다.")

    # 2
    slide = new_slide(prs)
    add_title(slide, "Why final-pose recognition is too late", "Embodied response requires early intention prediction plus actuator feasibility.", speaker="MinKyu")
    add_metric(slide, "response-window decision latency", "0.033 s", 0.75, 1.55, 2.55, accent=BLUE)
    add_metric(slide, "remaining actuator budget", "0.467 s", 3.55, 1.55, 2.55, accent=GREEN)
    add_metric(slide, "required response time", "0.413 s", 6.35, 1.55, 2.55, accent=AMBER)
    add_metric(slide, "hard deadline", "0.500 s", 9.15, 1.55, 2.55, accent=VIOLET)
    add_text(slide, "The robot must commit before the human motion is fully complete.", 1.1, 3.05, 11.0, 0.45, size=25, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["Sparse real videos make temporal learning difficult.", "The decision is only useful if the robot can still move.", "Therefore the system evaluates prediction, confidence, and actuator feasibility together."], 1.3, 3.88, 10.8, 1.35, size=20)
    add_image_panel(slide, FIGURES / "final_take_timing_budget.png", 2.0, 5.35, 9.3, 1.25)
    add_footer(slide, 2)
    set_notes(slide, "MinKyu Cho", "핵심 문제는 마지막 자세를 인식하는 것만으로는 늦다는 점입니다. 모델이 맞게 예측하더라도 로봇 손이 제한 시간 안에 움직이지 못하면 embodied counterattack이라고 보기 어렵습니다. 그래서 저희는 예측 정확도뿐 아니라 남은 시간, delay, 필요한 joint motion까지 같이 봤습니다.")

    # 3
    slide = new_slide(prs)
    add_title(slide, "End-to-end pipeline", "Few real clips become a live actuator-constrained counterattack system.", speaker="SeungHwan")
    add_image_panel(slide, FIGURES / "verified_system_pipeline.png", 0.72, 1.18, 11.9, 5.55)
    add_footer(slide, 3)
    set_notes(slide, "SeungHwan Kim", "전체 흐름은 적은 수의 실제 hand video에서 시작합니다. MediaPipe로 skeleton을 추출하고, canonical feature로 변환한 뒤, real-guided augmentation을 통해 synthetic skeleton sequence를 늘립니다. 그 다음 SCHUNK/Isaac-style pose progress에 맞춰 alignment하고, 최종적으로 TCN predictor와 robot counterattack policy로 연결합니다.")

    # 4
    slide = new_slide(prs)
    add_title(slide, "Real seed capture: small but verified", "The project starts from a deliberately small real-video seed set.", speaker="SeungHwan")
    add_metric(slide, "reviewed real clips", "20", 0.8, 1.45, 2.4, accent=BLUE)
    add_metric(slide, "processed frames", "720", 3.45, 1.45, 2.4, accent=TEAL)
    add_metric(slide, "detection coverage", "1.0", 6.1, 1.45, 2.4, accent=GREEN)
    add_metric(slide, "seed transitions", "2", 8.75, 1.45, 2.4, accent=VIOLET)
    add_bullets(slide, ["Real clips cover rock-to-paper and rock-to-scissors transitions.", "The seed set is intentionally few-shot, so the augmentation stage is central.", "Heldout test MP4s stay validation-only."], 0.95, 3.0, 5.2, 1.8, size=20)
    add_image_panel(slide, FIGURES / "canonical_skeleton_seed.png", 6.55, 2.75, 5.75, 3.15, label="Canonical seed preview from real MediaPipe landmarks")
    add_footer(slide, 4)
    set_notes(slide, "SeungHwan Kim", "여기서 중요한 것은 데이터가 많지 않았다는 점입니다. 실제 촬영은 20개 clip, 720 frame 수준이었고, 대신 각 frame에서 MediaPipe hand landmark가 안정적으로 잡혔는지를 먼저 검증했습니다. 이 작은 seed가 이후 augmentation의 기준이 됩니다.")

    # 5
    slide = new_slide(prs)
    add_title(slide, "MediaPipe skeleton representation", "Live video is converted into canonical temporal features, not raw image labels.", speaker="SeungHwan")
    add_image_panel(slide, FIGURES / "canonical_skeleton_seed.png", 0.78, 1.28, 6.05, 4.35)
    add_panel(slide, 7.2, 1.42, 4.85, 3.9, fill="ECFEFF", line="99F6E4")
    add_text(slide, "Feature vector", 7.52, 1.72, 2.8, 0.34, size=22, bold=True, font_color=TEAL)
    add_bullets(slide, ["21 canonical landmarks", "per-frame velocities", "five finger-curl values", "tip-to-MCP distances", "fingertip speeds", "motion signal"], 7.55, 2.22, 4.1, 2.35, size=18)
    add_metric(slide, "per-frame dimensions", "142-D", 7.5, 5.55, 2.2, accent=TEAL)
    add_metric(slide, "canonical sequences", "20", 10.0, 5.55, 1.7, accent=BLUE)
    add_footer(slide, 5)
    set_notes(slide, "SeungHwan Kim", "MediaPipe에서 얻은 21개 landmark를 그대로 쓰는 것이 아니라 wrist 기준 좌표계로 정규화했습니다. 또한 velocity, finger curl, fingertip speed 같은 temporal feature를 붙여서 한 frame당 142차원 feature를 만들었습니다. 이 구조 덕분에 카메라 위치나 크기 변화에 조금 더 강한 입력을 만들 수 있었습니다.")

    # 6
    slide = new_slide(prs)
    add_title(slide, "Real-guided skeleton augmentation", "The few-shot seed set is expanded into trainable temporal variation.", speaker="SeungHwan")
    add_image_panel(slide, FIGURES / "real_guided_skeleton_augmentation.png", 0.75, 1.25, 5.55, 4.75, label="Bulk augmented skeleton trajectories")
    add_metric(slide, "compact samples", "2,000", 6.75, 1.55, 2.55, accent=BLUE)
    add_metric(slide, "large sharded samples", "10,000", 9.6, 1.55, 2.75, accent=VIOLET)
    add_bullets(slide, ["Temporal scaling and warping preserve real motion structure.", "Finger-specific scaling creates controlled gesture variation.", "Noise and small rotations improve robustness without changing labels.", "Large split: 7,000 train / 1,500 validation / 1,500 test."], 6.8, 2.95, 5.25, 2.35, size=18)
    add_footer(slide, 6)
    set_notes(slide, "SeungHwan Kim", "이 단계가 연구의 핵심입니다. 저희는 임의로 손 그림을 생성한 것이 아니라 실제 skeleton trajectory를 seed로 두고 temporal scaling, finger-specific scaling, noise를 적용했습니다. 그래서 적은 실제 데이터에서 출발하지만 motion pattern은 실제 영상에 기반한 synthetic dataset을 만들 수 있었습니다.")

    # 7
    slide = new_slide(prs)
    add_title(slide, "Isaac/SCHUNK alignment makes the data render-ready", "Skeleton progress is aligned to a robot-hand response schedule.", speaker="SeungHwan")
    add_image_panel(slide, FIGURES / "real_guided_alignment_error.png", 0.75, 1.28, 5.55, 3.1, label="Progress-error comparison")
    add_image_panel(slide, FIGURES / "real_guided_pose_schedule.png", 6.65, 1.28, 5.55, 3.1, label="Real-guided pose schedule")
    add_metric(slide, "old max progress error", "~0.106", 1.3, 4.85, 2.75, accent=RED)
    add_metric(slide, "32-frame max error", "~0.016", 4.45, 4.85, 2.75, accent=GREEN)
    add_metric(slide, "render-manifest entries", "64", 7.6, 4.85, 2.35, accent=BLUE)
    add_text(slide, "This stage connects skeleton learning to the later SCHUNK/Isaac response layer without claiming fresh image-conditioned training.", 1.0, 6.17, 11.3, 0.38, size=17, font_color=SLATE, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)
    set_notes(slide, "SeungHwan Kim", "증강된 skeleton을 로봇 손과 연결하기 위해 progress schedule도 맞췄습니다. 이전 8-frame grid에서는 progress error가 컸지만 32-frame real-guided schedule로 줄였습니다. 다만 이 단계는 이미지 기반 Isaac training을 완료했다는 뜻은 아니고, skeleton dataset이 향후 render-conditioned training으로 이어질 수 있게 준비된 상태라는 의미입니다.")

    # 8
    slide = new_slide(prs)
    add_title(slide, "Temporal predictor and live policy", "The final demo uses frozen v4; v7e stays diagnostic-only.", speaker="MinKyu")
    add_image_panel(slide, FIGURES / "tcn_temporal_model.png", 0.82, 1.25, 4.35, 4.5, label="TCN-style temporal model schematic")
    add_panel(slide, 5.65, 1.3, 6.55, 4.45, fill="F5F3FF", line="DDD6FE")
    add_text(slide, "Frozen live/demo policy", 5.98, 1.62, 4.3, 0.36, size=22, bold=True, font_color=VIOLET)
    add_bullets(slide, ["v4 fallback predictor", "profile weights: 0.25 / 0.75 / 0.0", "confidence >= 0.70, margin >= 0.10", "confirmation count = 2", "response prompt = scissors", "v7e: diagnostics only, original20 = 17/20"], 6.0, 2.12, 5.7, 2.65, size=18)
    add_footer(slide, 8)
    set_notes(slide, "MinKyu Cho", "최종 live/demo 모델은 v4 fallback policy로 고정했습니다. v7e는 original20에서 17/20으로 strict gate를 통과하지 못했기 때문에 report diagnostic으로만 남겼고, v7f retraining은 시작하지 않았습니다. 발표에서는 이 경계를 명확히 말하는 것이 중요합니다.")

    # 9
    slide = new_slide(prs)
    add_title(slide, "Counterattack policy", "The robot starts from rock and only moves after a confirmed response-window decision.", speaker="MinKyu")
    rows = [
        ["Human prediction", "Robot counter", "Outcome"],
        ["rock", "paper", "robot wins"],
        ["paper", "scissors", "robot wins"],
        ["scissors", "rock", "robot wins"],
    ]
    add_simple_table(slide, rows, 0.9, 1.45, 5.25, 2.55)
    add_image_panel(slide, FIGURES / "schunk_rock_yaw45_pitch20.png", 6.55, 1.3, 1.85, 2.2, label="rock")
    add_arrow(slide, 8.5, 2.05, 0.55, 0.32, fill=SLATE)
    add_image_panel(slide, FIGURES / "schunk_paper_yaw45_pitch20.png", 9.12, 1.3, 1.85, 2.2, label="paper")
    add_arrow(slide, 8.5, 4.65, 0.55, 0.32, fill=SLATE)
    add_image_panel(slide, FIGURES / "schunk_scissors_yaw45_pitch20.png", 9.12, 3.9, 1.85, 2.2, label="scissors")
    add_text(slide, "Wait states map to robot rock; no counter move is issued until a response-window prediction is confirmed.", 0.95, 4.75, 5.15, 0.78, size=21, bold=True, font_color=INK)
    add_footer(slide, 9)
    set_notes(slide, "MinKyu Cho", "counterattack policy는 간단하지만 timing gate와 함께 쓰일 때 의미가 있습니다. 사람의 예측이 rock이면 robot은 paper, paper면 scissors, scissors면 rock을 냅니다. 대기 상태에서는 항상 robot rock으로 시작하고, response window에서 확정된 decision이 있을 때만 움직입니다.")

    # 10
    slide = new_slide(prs)
    add_title(slide, "Actuator timing model", "A correct prediction is accepted only if the robot can still move in time.", speaker="MinKyu")
    add_image_panel(slide, FIGURES / "final_take_timing_budget.png", 0.9, 1.25, 7.6, 3.25)
    add_panel(slide, 8.85, 1.33, 3.45, 3.05, fill="FFF7ED", line="FDBA74")
    add_text(slide, "Feasibility rule", 9.15, 1.62, 2.4, 0.32, size=22, bold=True, font_color=AMBER)
    add_bullets(slide, ["response delay = 0.08 s", "deadline = 0.50 s", "required time = 0.413 s", "remaining time = 0.467 s", "limiting joint = index_curl"], 9.16, 2.15, 2.75, 1.58, size=16)
    add_text(slide, "Measured timing is local desktop end-to-end timing: camera + MediaPipe + features + PyTorch/CUDA + output.", 1.05, 5.15, 10.9, 0.68, size=22, bold=True, font_color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)
    set_notes(slide, "MinKyu Cho", "여기서는 모델 forward time만 말하는 것이 아니라 local desktop loop 전체를 말합니다. 카메라 입력, MediaPipe, feature construction, PyTorch inference, OpenCV output이 모두 포함됩니다. 그래서 이 숫자는 hardware-independent latency가 아니라 이 PC에서 검증한 end-to-end timing입니다.")

    # 11
    slide = new_slide(prs)
    add_title(slide, "Final live demo: human rock -> robot paper", "One selected live take passed the full actuator-feasible counterattack gate.", speaker="MinKyu")
    add_image_panel(slide, FIGURES / "final_live_candidate_poster.png", 0.82, 1.22, 7.2, 4.75, label="Final selected live take poster frame")
    add_panel(slide, 8.42, 1.35, 3.85, 3.75, fill="ECFDF5", line="BBF7D0")
    add_text(slide, "Selected take", 8.75, 1.65, 2.4, 0.35, size=22, bold=True, font_color=GREEN)
    add_bullets(slide, ["human target: rock", "robot counter: paper", "prediction: rock", "confidence: 0.77", "remaining time: 0.467 s", "result: actuator-feasible win"], 8.78, 2.12, 3.0, 2.0, size=17)
    add_text(slide, "During recording, play or cut in the final demo video after this slide.", 8.6, 5.38, 3.5, 0.35, size=14, font_color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)
    set_notes(slide, "MinKyu Cho", "이 슬라이드에서는 최종 demo video를 같이 보여주면 좋습니다. 선택된 take는 prompt scissors response window에서 제가 rock을 냈고, 모델이 rock으로 확정한 뒤 robot paper로 counterattack 한 사례입니다. 이 take는 feasibility check까지 통과했습니다.")

    # 12
    slide = new_slide(prs)
    add_title(slide, "Results snapshot", "Dataset scale, policy boundaries, and live outcomes are reported together.", speaker="MinKyu")
    rows = [
        ["Evidence", "Value", "Role"],
        ["Real clips / frames", "20 / 720", "seed capture"],
        ["Feature dimension", "142-D", "predictor input"],
        ["Large dataset", "10,000", "real-guided training base"],
        ["v4 policy", "live/demo", "frozen fallback"],
        ["v7e", "17/20", "diagnostics only"],
        ["selected take", "actuator-feasible win", "final demo"],
    ]
    add_simple_table(slide, rows, 0.75, 1.25, 6.2, 4.55, header_fill="EDE9FE")
    add_image_panel(slide, FIGURES / "three_take_review_poster.png", 7.35, 1.35, 4.7, 2.6, label="Three-take review poster")
    add_bullets(slide, ["Paper take: no confirmed paper response-window decision.", "Scissors take: confirmed scissors, but too late for the actuator deadline.", "Failures are kept as diagnostics, not hidden."], 7.42, 4.35, 4.55, 1.28, size=17)
    add_footer(slide, 12)
    set_notes(slide, "MinKyu Cho", "결과는 성공 사례만 말하지 않고 전체 evidence를 같이 보여주는 것이 좋습니다. v4는 최종 live/demo 모델이고, v7e는 17/20으로 diagnostic branch입니다. 세 live take 중 최종 선택된 것은 actuator-feasible win이지만, paper와 scissors 실패도 timing과 response-window 문제가 무엇인지 보여주는 진단 자료입니다.")

    # 13
    slide = new_slide(prs)
    add_title(slide, "Limitations", "The project is honest about what was validated and what remains open.", speaker="MinKyu")
    add_panel(slide, 0.85, 1.35, 3.7, 3.75, fill="FEF2F2", line="FECACA")
    add_text(slide, "Skeleton-first", 1.15, 1.72, 2.5, 0.35, size=24, bold=True, font_color=RED)
    add_bullets(slide, ["The final predictor does not learn from raw Isaac camera images.", "Render-ready sidecars prepare future image-conditioned work."], 1.16, 2.35, 2.95, 1.55, size=18)
    add_panel(slide, 4.85, 1.35, 3.7, 3.75, fill="EFF6FF", line="BFDBFE")
    add_text(slide, "Local timing", 5.15, 1.72, 2.5, 0.35, size=24, bold=True, font_color=BLUE)
    add_bullets(slide, ["Latency depends on camera, MediaPipe, OpenCV, PyTorch/CUDA, and PC specs.", "It is not a universal model-speed claim."], 5.16, 2.35, 2.95, 1.55, size=18)
    add_panel(slide, 8.85, 1.35, 3.7, 3.75, fill="FFFBEB", line="FDE68A")
    add_text(slide, "Live breadth", 9.15, 1.72, 2.5, 0.35, size=24, bold=True, font_color=AMBER)
    add_bullets(slide, ["Only one final selected take passed the full gate.", "More subjects and more controlled takes are needed."], 9.16, 2.35, 2.95, 1.55, size=18)
    add_text(slide, "Boundary: no v7f retraining and no v7e promotion in this submission.", 1.2, 5.75, 10.9, 0.44, size=23, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 13)
    set_notes(slide, "MinKyu Cho", "한계를 명확히 말하겠습니다. 첫째, 최종 predictor는 skeleton-first입니다. 둘째, timing은 제 local desktop에서의 end-to-end timing입니다. 셋째, live take는 하나만 full gate를 통과했습니다. 그래서 다음 단계는 더 많은 subject와 hardware-normalized latency benchmark가 필요합니다.")

    # 14
    slide = new_slide(prs)
    set_bg(slide, "0F172A")
    add_text(slide, "What this project proves", 0.78, 0.75, 6.0, 0.5, size=32, bold=True, font_color="F8FAFC")
    add_text(slide, "A small verified real skeleton seed set can be expanded, aligned to a robot-hand embodiment, and transferred back to live prompt-window input for an actuator-feasible counterattack.", 0.82, 1.45, 6.2, 1.15, size=21, font_color="E2E8F0")
    add_panel(slide, 7.35, 0.88, 4.95, 4.6, fill="111827", line="334155")
    add_text(slide, "Next steps", 7.7, 1.23, 2.8, 0.35, size=24, bold=True, font_color="F8FAFC")
    add_bullets(slide, ["Hardware-normalized latency benchmark", "Image-conditioned Isaac/SCHUNK training", "Physical hand deployment with calibrated actuators", "Broader live-subject validation"], 7.7, 1.95, 3.95, 2.35, size=19, font_color="E2E8F0")
    add_chip(slide, "Dataset links: fill after Google Drive upload", 0.85, 4.0, 3.8, 0.4, fill="DBEAFE", font_color="1D4ED8")
    add_chip(slide, "Demo/presentation links: fill after YouTube upload", 0.85, 4.55, 4.35, 0.4, fill="DCFCE7", font_color="166534")
    add_text(slide, "Thank you", 0.82, 5.7, 3.0, 0.55, size=34, bold=True, font_color="F8FAFC")
    add_footer(slide, 14)
    set_notes(slide, "MinKyu Cho", "마무리하겠습니다. 이 프로젝트의 핵심은 적은 실제 skeleton seed에서 출발해 real-guided synthetic dataset을 만들고, 이것을 robot-hand timing constraint와 연결한 점입니다. 앞으로는 latency를 더 세밀하게 분리해서 측정하고, image-conditioned Isaac training과 실제 로봇 손 deployment로 확장할 수 있습니다. 감사합니다.")

    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    build()
